"""
VerbaFlow AI - Document Processor Service
Extracts text, metadata, and validates files for all supported formats.
"""
from __future__ import annotations

import csv
import hashlib
import io
import logging
import os
import re
import struct
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

logger = logging.getLogger(__name__)

# File magic signatures for safety scanning
MAGIC_SIGNATURES = {
    "pdf": b"%PDF",
    "docx": b"PK\x03\x04",  # ZIP-based Office docs
    "xlsx": b"PK\x03\x04",
    "pptx": b"PK\x03\x04",
    "zip": b"PK\x03\x04",
}

DANGEROUS_EXTENSIONS = {".exe", ".bat", ".sh", ".ps1", ".dll", ".msi", ".vbs"}


class DocumentProcessor:
    """
    Extracts text and metadata from supported file formats.

    Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, HTML, Markdown.
    """

    async def extract_text(self, file_path: str, file_type: str) -> str:
        """
        Route extraction to the appropriate handler based on file_type.

        Args:
            file_path: Absolute path to the file.
            file_type: One of: pdf, docx, txt, csv, xlsx, pptx, html, md.

        Returns:
            Extracted plain text string.
        """
        extractors = {
            "pdf": self.extract_text_from_pdf,
            "docx": self.extract_text_from_docx,
            "txt": self.extract_text_from_txt,
            "csv": self.extract_text_from_csv,
            "xlsx": self.extract_text_from_xlsx,
            "pptx": self.extract_text_from_pptx,
            "html": self.extract_text_from_html,
            "md": self.extract_text_from_markdown,
        }
        extractor = extractors.get(file_type.lower())
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        return await extractor(file_path)

    async def extract_text_from_pdf(self, path: str) -> str:
        """
        Extract all text from a PDF file using pypdf.

        Args:
            path: Path to the PDF file.

        Returns:
            Concatenated text from all pages.
        """
        try:
            from pypdf import PdfReader

            reader = PdfReader(path)
            pages: list[str] = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[Page {i + 1}]\n{text}")
            return "\n\n".join(pages)
        except Exception as exc:
            logger.error("PDF extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"PDF extraction failed: {exc}") from exc

    async def extract_text_from_docx(self, path: str) -> str:
        """
        Extract text from a DOCX file using python-docx.

        Preserves paragraph structure and table content.

        Args:
            path: Path to the DOCX file.

        Returns:
            Extracted text with paragraph breaks.
        """
        try:
            from docx import Document

            doc = Document(path)
            parts: list[str] = []

            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            # Extract table cells
            for table in doc.tables:
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        parts.append(" | ".join(row_texts))

            return "\n\n".join(parts)
        except Exception as exc:
            logger.error("DOCX extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"DOCX extraction failed: {exc}") from exc

    async def extract_text_from_txt(self, path: str) -> str:
        """
        Read a plain text file, trying UTF-8 then Latin-1 encoding.

        Args:
            path: Path to the text file.

        Returns:
            File content as a string.
        """
        try:
            try:
                with open(path, "r", encoding="utf-8") as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(path, "r", encoding="latin-1") as f:
                    return f.read()
        except Exception as exc:
            logger.error("TXT extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"TXT extraction failed: {exc}") from exc

    async def extract_text_from_csv(self, path: str) -> str:
        """
        Convert a CSV file to a readable text representation.

        Includes the header row and all data rows. Large files are truncated
        at 10,000 rows to avoid memory issues.

        Args:
            path: Path to the CSV file.

        Returns:
            Tab-separated text with header and data rows.
        """
        try:
            import pandas as pd

            df = pd.read_csv(path, nrows=10000)
            lines = [" | ".join(str(v) for v in df.columns)]
            for _, row in df.iterrows():
                lines.append(" | ".join(str(v) for v in row.values))
            return "\n".join(lines)
        except Exception as exc:
            logger.error("CSV extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"CSV extraction failed: {exc}") from exc

    async def extract_text_from_xlsx(self, path: str) -> str:
        """
        Extract text from all sheets of an XLSX file.

        Args:
            path: Path to the XLSX file.

        Returns:
            Text content from all sheets.
        """
        try:
            import openpyxl

            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            all_text: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text.append(f"=== Sheet: {sheet_name} ===")
                for row in ws.iter_rows(values_only=True):
                    row_values = [str(c) for c in row if c is not None]
                    if row_values:
                        all_text.append(" | ".join(row_values))
            return "\n".join(all_text)
        except Exception as exc:
            logger.error("XLSX extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"XLSX extraction failed: {exc}") from exc

    async def extract_text_from_pptx(self, path: str) -> str:
        """
        Extract text from all slides and notes of a PPTX file.

        Args:
            path: Path to the PPTX file.

        Returns:
            Text content with slide markers.
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt

            prs = Presentation(path)
            slides_text: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content: list[str] = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs).strip()
                            if text:
                                slide_content.append(text)
                    # Tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_content.append(" | ".join(cells))
                # Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[Notes] {notes}")
                slides_text.append("\n".join(slide_content))
            return "\n\n".join(slides_text)
        except Exception as exc:
            logger.error("PPTX extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"PPTX extraction failed: {exc}") from exc

    async def extract_text_from_html(self, path: str) -> str:
        """
        Parse an HTML file and extract visible text using BeautifulSoup.

        Removes script, style, and navigation elements.

        Args:
            path: Path to the HTML file.

        Returns:
            Visible text content.
        """
        try:
            from bs4 import BeautifulSoup

            with open(path, "r", encoding="utf-8", errors="replace") as f:
                html = f.read()

            soup = BeautifulSoup(html, "html.parser")
            # Remove non-content elements
            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)
            # Collapse multiple blank lines
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text.strip()
        except Exception as exc:
            logger.error("HTML extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"HTML extraction failed: {exc}") from exc

    async def extract_text_from_markdown(self, path: str) -> str:
        """
        Convert Markdown to plain text by stripping MD syntax.

        Args:
            path: Path to the Markdown file.

        Returns:
            Plain text with Markdown formatting removed.
        """
        try:
            import markdown
            from bs4 import BeautifulSoup

            with open(path, "r", encoding="utf-8") as f:
                md_text = f.read()

            html = markdown.markdown(md_text, extensions=["tables", "fenced_code"])
            soup = BeautifulSoup(html, "html.parser")
            return soup.get_text(separator="\n", strip=True)
        except Exception as exc:
            logger.error("Markdown extraction failed for %s: %s", path, exc)
            raise RuntimeError(f"Markdown extraction failed: {exc}") from exc

    async def extract_metadata(
        self, path: str, file_type: str
    ) -> Dict[str, Any]:
        """
        Extract file metadata: title, author, page count, creation date, etc.

        Args:
            path: File path.
            file_type: File type string.

        Returns:
            Dict with metadata keys.
        """
        metadata: Dict[str, Any] = {
            "file_size": os.path.getsize(path),
            "file_type": file_type,
            "sha256": self._compute_sha256(path),
        }

        try:
            if file_type == "pdf":
                from pypdf import PdfReader
                reader = PdfReader(path)
                info = reader.metadata
                metadata["page_count"] = len(reader.pages)
                if info:
                    metadata["title"] = info.get("/Title", "")
                    metadata["author"] = info.get("/Author", "")
                    metadata["creator"] = info.get("/Creator", "")
                    creation = info.get("/CreationDate", "")
                    metadata["creation_date"] = str(creation)

            elif file_type == "docx":
                from docx import Document
                doc = Document(path)
                props = doc.core_properties
                metadata["title"] = props.title or ""
                metadata["author"] = props.author or ""
                metadata["created"] = str(props.created) if props.created else ""
                metadata["paragraph_count"] = len(doc.paragraphs)

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(path)
                metadata["slide_count"] = len(prs.slides)
                props = prs.core_properties
                metadata["title"] = props.title or ""
                metadata["author"] = props.author or ""

            elif file_type == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True)
                metadata["sheet_count"] = len(wb.sheetnames)
                metadata["sheets"] = wb.sheetnames

        except Exception as exc:
            logger.warning("Metadata extraction partial failure: %s", exc)

        return metadata

    def clean_text(self, text: str) -> str:
        """
        Normalise extracted text by removing noise.

        Operations:
        - Decode HTML entities
        - Collapse multiple whitespace characters
        - Remove null bytes and control characters
        - Strip leading/trailing whitespace

        Args:
            text: Raw extracted text.

        Returns:
            Cleaned text string.
        """
        import html

        # Decode HTML entities
        text = html.unescape(text)
        # Remove null bytes
        text = text.replace("\x00", "")
        # Remove other control characters (keep newlines and tabs)
        text = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        # Collapse multiple spaces on the same line
        text = re.sub(r"[ \t]+", " ", text)
        # Collapse more than 2 consecutive newlines into 2
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def scan_file_safety(self, path: str) -> bool:
        """
        Basic file safety check using:
        1. Extension blocklist for clearly dangerous extensions.
        2. Magic byte header validation for binary formats.

        Args:
            path: Absolute path to the file.

        Returns:
            True if safe, False if suspicious.
        """
        p = Path(path)

        # Check extension blocklist
        if p.suffix.lower() in DANGEROUS_EXTENSIONS:
            logger.warning("Dangerous file extension blocked: %s", path)
            return False

        # Verify magic bytes for binary formats
        try:
            with open(path, "rb") as f:
                header = f.read(8)

            ext = p.suffix.lower().lstrip(".")
            if ext in MAGIC_SIGNATURES:
                expected = MAGIC_SIGNATURES[ext]
                if not header.startswith(expected):
                    logger.warning(
                        "Magic byte mismatch for %s: expected %s, got %s",
                        path,
                        expected.hex(),
                        header[:4].hex(),
                    )
                    return False
        except (OSError, IOError) as exc:
            logger.error("File safety scan IO error: %s", exc)
            return False

        # Check file size (100MB hard limit for safety scan)
        file_size = os.path.getsize(path)
        if file_size > 100 * 1024 * 1024:
            logger.warning("File exceeds maximum safe size: %s (%d bytes)", path, file_size)
            return False

        return True

    def _compute_sha256(self, path: str) -> str:
        """Compute SHA-256 hash of a file for deduplication."""
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()

    async def count_words(self, text: str) -> int:
        """Estimate word count from extracted text."""
        return len(text.split())

    async def detect_language(self, text: str) -> str:
        """
        Detect the primary language of the text.
        Returns ISO 639-1 code or 'en' as fallback.
        """
        try:
            # Use a simple heuristic based on common English words
            # In production, integrate langdetect or fasttext
            sample = text[:500].lower()
            english_words = {"the", "is", "are", "and", "or", "of", "in", "to", "a"}
            word_set = set(re.findall(r"\b\w+\b", sample))
            if len(english_words & word_set) > 3:
                return "en"
            return "en"  # Default fallback
        except Exception:
            return "en"
